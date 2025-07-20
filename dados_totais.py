# %%
# Importacao das bibliotecas

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import os
import time

# %%
# lendo o dataframe

df = pd.read_excel("data_frame_total.xlsx")
# df.head()


# Lista de palavras a ignorar como "segundo nome"
palavras_ignoradas = ['da', 'de', 'do', 'das', 'dos']

# Exemplo de DataFrame

# Função para pegar o primeiro e segundo nome (considerando palavras ignoradas)


def extrair_primeiros_nomes(nome):
    partes = nome.split()
    if len(partes) > 2 and partes[1].lower() in palavras_ignoradas:
        return ' '.join(partes[:3])  # Pega o primeiro, segundo e terceiro
    else:
        return ' '.join(partes[:2])  # Pega apenas o primeiro e o segundo


# Aplicar a função na coluna de nomes
df['nome'] = df['nome'].apply(extrair_primeiros_nomes)

# %%
# Gerando os graficos

sns.set_style("whitegrid")

# Grafico 1 : Notas por aluno por bimestre


def gerar_graficos_turma(turma: str):
    df_turma = df[df['turma'] == int(turma)]

    # Pivotando a tabela para ter cada aluno em uma linha e os bimestres em colunas
    df_pivot = df_turma.pivot(index='nome', columns='bimestre',
                              values='Notas').sort_values(by='nome', ascending=False)

    fig, ax = plt.subplots(1, 1)
    df_pivot.plot(kind='barh', stacked=True, figsize=(
        10, 6), colormap='viridis', ax=ax)
    ax.set_title(f"Turma {str(turma)} - Notas Bimestrais")
    ax.set_ylabel("")
    ax.set_xlabel("Nota")
    # Save the chart
    plt.tight_layout()
    output_path = "imagens_totais"
    output_file = f"{output_path}/notas_{turma}.png"
    fig.savefig(output_file, format="png", dpi=500)
    plt.close(fig)
    print(f"Gráfico das notas da turma {turma} salvo em: {output_file}")

    # Grafico 2: Grafico de faltas bimestrais por aluno

    df_faltas = df_turma.pivot(index='nome', columns='bimestre',
                               values='Faltas').sort_values(by='nome', ascending=False)

    fig2, ax2 = plt.subplots(1, 1)
    df_faltas.plot(kind='barh', stacked=True, figsize=(
        10, 6), colormap='viridis', ax=ax2)
    ax2.set_title(f"Turma {str(turma)} - Faltas Bimestrais")
    ax2.set_ylabel("")
    ax2.set_xlabel("Faltas")
    # Save the chart
    plt.tight_layout()
    output_path = "imagens_totais"
    output_file = f"{output_path}/faltas_{turma}.png"
    fig2.savefig(output_file, format="png", dpi=500)
    plt.close(fig2)
    print(f"Gráfico das faltas da turma {turma} salvo em: {output_file}")


def gerar_graficos_media(df):
    # Grafico 3: Média de nota por turma por bimestre
    df_media = df.groupby(['bimestre', 'turma'])['Notas'].mean().reset_index()
    fig3, ax3 = plt.subplots(1, 1)
    sns.barplot(x='bimestre', y='Notas', hue='turma',
                data=df_media, ax=ax3, palette='viridis')
    ax3.set_title(f"Média de Notas por Turma")
    ax3.set_ylabel("Média de Notas")
    ax3.set_xlabel("Bimestre")

    plt.tight_layout()
    output_path = "imagens_totais"
    output_file = f"{output_path}/medias_notas_bimestre.png"
    fig3.savefig(output_file, format="png", dpi=500)
    plt.close(fig3)
    print(f"Gráfico das medias de notas foi salvo em: {output_file}")

    # Grafico 4: Média de falta por turma por bimestre
    df_media = df.groupby(['bimestre', 'turma'])['Faltas'].mean().reset_index()
    fig3, ax3 = plt.subplots(1, 1)
    sns.barplot(x='bimestre', y='Faltas', hue='turma',
                data=df_media, ax=ax3, palette='viridis')
    ax3.set_title(f"Média de Faltas por Turma")
    ax3.set_ylabel("Média de Faltas")
    ax3.set_xlabel("Bimestre")

    plt.tight_layout()
    output_path = "imagens_totais"
    output_file = f"{output_path}/medias_faltas_bimestre.png"
    fig3.savefig(output_file, format="png", dpi=500)
    plt.close(fig3)
    print(f"Gráfico das medias de faltas foi salvo em: {output_file}")


# %%
gerar_graficos_turma('900')
gerar_graficos_turma('901')
gerar_graficos_turma('902')
gerar_graficos_media(df)


def verificar_arquivo_aberto(caminho_arquivo):
    try:
        # Tentar abrir o arquivo exclusivamente
        with open(caminho_arquivo, 'r+') as arquivo:
            try:
                # Tentativa de renomear o arquivo para verificar se está aberto
                os.rename(caminho_arquivo, caminho_arquivo)
                print("O arquivo não está aberto. Você pode prosseguir.")
                return False
            except OSError:
                print("Erro: O arquivo já está aberto por outro processo.")
                return True
    except IOError:
        print(f"Erro: Não foi possível abrir o arquivo {caminho_arquivo}.")
        return True


# %%
# Gerando o power point


# Criar a apresentação
prs = Presentation()

# Configurar o tamanho do slide para widescreen (16:9)
prs.slide_width = Inches(13.33)  # Largura em widescreen
prs.slide_height = Inches(7.5)  # Altura em widescreen

# Slide 1: Slide de título
slide_layout = prs.slide_layouts[6]  # Layout em branco (para controle total)
slide = prs.slides.add_slide(slide_layout)

# Adicionar título
title_box = slide.shapes.add_textbox(
    Inches(1), Inches(2), Inches(11.33), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Gráficos de Notas e Faltas - 9º ano"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centralizar texto
title_frame.paragraphs[0].font.size = Pt(60)  # Ajustar tamanho da fonte

# Adicionar subtítulo
subtitle_box = slide.shapes.add_textbox(
    Inches(1), Inches(4), Inches(11.33), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "4º bimestre - E. M. A. Pereira Bruno"
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centralizar texto
subtitle_frame.paragraphs[0].font.size = Pt(44)  # Ajustar tamanho da fonte

# Imagens (Substitua os caminhos abaixo pelos arquivos das suas imagens)
imagens = [
    "imagens_totais/notas_900.png",
    "imagens_totais/notas_901.png",
    "imagens_totais/notas_902.png",
    "imagens_totais/faltas_900.png",
    "imagens_totais/faltas_901.png",
    "imagens_totais/faltas_902.png",
    "imagens_totais/medias_notas_bimestre.png",
    "imagens_totais/medias_faltas_bimestre.png"
]

# Adicionar os slides com imagens
for img_path in imagens:
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)

    # Definir a posição e o tamanho da imagem para cobrir todo o slide
    left = 0  # Margem esquerda
    top = 0   # Margem superior
    slide.shapes.add_picture(
        img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# Salvar a apresentação

prs.save("apresentacao_totais.pptx")

# Abrir o arquivo automaticamente
file_path = "apresentacao_totais.pptx"
print("Apresentação criada com sucesso!")
print("Abrindo apresentação, aguarde...")
time.sleep(2)
os.startfile(file_path)
